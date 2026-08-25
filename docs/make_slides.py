"""Build the 10-slide presentation deck (no title slide - author supplies it).

Widescreen 16:9, designed to be readable from the back of a room: one idea
per slide, large type, and every number traceable to the report.

    python docs/make_slides.py
"""
import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

HERE = pathlib.Path(__file__).resolve().parent
FIG = HERE / "ieee_report" / "figures"
OUT = HERE / "chess_ai_slides.pptx"

# ---- palette ---------------------------------------------------------
INK = RGBColor(0x1A, 0x1D, 0x24)      # near-black body text
MUTED = RGBColor(0x5A, 0x64, 0x72)    # secondary text
BLUE = RGBColor(0x1F, 0x4E, 0x9C)     # headings / primary accent
GREEN = RGBColor(0x1A, 0x7F, 0x4B)    # good outcomes
RED = RGBColor(0xB3, 0x35, 0x2C)      # failures
AMBER = RGBColor(0xB8, 0x7A, 0x14)
BG = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xF2, 0xF5, 0xF9)    # tinted callout panel

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.75)
BODY_W = W - 2 * MARGIN

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def textbox(s, left, top, width, height):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size, color=INK, bold=False, space_before=0,
         space_after=6, align=PP_ALIGN.LEFT, first=False, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Segoe UI"
    return p


def header(s, kicker, title, title_color=INK):
    """Small coloured kicker above a large headline, plus a rule."""
    tf = textbox(s, MARGIN, Inches(0.45), BODY_W, Inches(1.35))
    para(tf, kicker.upper(), 13, BLUE, bold=True, space_after=4, first=True)
    para(tf, title, 33, title_color, bold=True, space_after=0)
    line = s.shapes.add_shape(1, MARGIN, Inches(1.72), Inches(1.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    line.shadow.inherit = False


def bullets(s, items, top=Inches(2.1), left=MARGIN, width=None, size=19):
    """items: list of (text, color, bold) or plain strings."""
    tf = textbox(s, left, top, width or BODY_W, H - top - Inches(0.6))
    for i, it in enumerate(items):
        if isinstance(it, str):
            text, color, bold = it, INK, False
        else:
            text, color, bold = (list(it) + [INK, False])[:3]
        para(tf, text, size, color, bold=bold, space_after=13, first=(i == 0))
    return tf


def panel(s, left, top, width, height, color=PANEL):
    box = s.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    box.shadow.inherit = False
    return box


def stat(s, left, top, width, value, label, color=BLUE, vsize=44):
    """A large number with a caption underneath."""
    panel(s, left, top, width, Inches(1.65))
    tf = textbox(s, left, top + Inches(0.18), width, Inches(1.3))
    tf.word_wrap = True
    para(tf, value, vsize, color, bold=True, align=PP_ALIGN.CENTER,
         space_after=2, first=True)
    para(tf, label, 13, MUTED, align=PP_ALIGN.CENTER, space_after=0)


def picture(s, name, left, top, width):
    path = FIG / name
    if path.exists():
        s.shapes.add_picture(str(path), left, top, width=width)
    else:                      # keep the deck buildable if a figure is absent
        panel(s, left, top, width, Inches(3))
        tf = textbox(s, left, top + Inches(1.3), width, Inches(0.5))
        para(tf, f"[{name} not generated yet]", 14, MUTED,
             align=PP_ALIGN.CENTER, first=True)


def note(s, text, color=MUTED, size=15, italic=True):
    tf = textbox(s, MARGIN, H - Inches(1.0), BODY_W, Inches(0.55))
    para(tf, text, size, color, italic=italic, first=True)


def code(s, lines, left, top, width, size=17, color=INK):
    height = Inches(0.42) * len(lines) + Inches(0.36)
    panel(s, left, top, width, height, RGBColor(0xF7, 0xF7, 0xF4))
    tf = textbox(s, left + Inches(0.22), top + Inches(0.16),
                 width - Inches(0.4), height)
    for i, ln in enumerate(lines):
        text, c = (ln, color) if isinstance(ln, str) else ln
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.name = "Consolas"
        r.font.color.rgb = c


# =====================================================================
# 1 - What I built
# =====================================================================
s = slide()
header(s, "The project", "Teaching a network chess, with no chess knowledge")
bullets(s, [
    ("No opening book. No piece values. No evaluation function.", INK, True),
    "Input: the rules, plus human games. Everything else is learned.",
    "One network, two heads — policy (which moves look good) and value "
    "(who is winning).",
    "The network is not a player. Monte Carlo Tree Search turns it into one.",
    "Two stages: imitate humans, then improve by playing itself.",
], top=Inches(2.15), width=Inches(7.4))

stat(s, Inches(8.6), Inches(2.3), Inches(3.9), "760,717", "trainable parameters")
stat(s, Inches(8.6), Inches(4.15), Inches(3.9), "676,648",
     "human positions used for pretraining", GREEN)
note(s, "AlphaZero used ~5,000 TPUs. This runs on one desktop with an RTX 3060.")

# =====================================================================
# 2 - Representation
# =====================================================================
s = slide()
header(s, "Data science: representation", "Turning a chessboard into a tensor")
bullets(s, [
    ("Input — 18 binary planes of 8×8", BLUE, True),
    "12 piece planes (6 types × 2 colours)  ·  4 castling rights  "
    "·  1 en passant  ·  1 side to move",
    ("Output — 4,672 actions", BLUE, True),
    "64 origin squares × 73 move types: 56 queen-style, 8 knight, "
    "9 underpromotions.",
    "Every legal chess move maps to exactly one index. Verified by a "
    "round-trip test over thousands of positions: 0 collisions.",
], top=Inches(2.1), width=Inches(7.5), size=18)

stat(s, Inches(8.8), Inches(2.4), Inches(3.7), "18 × 8 × 8", "board in", BLUE, 36)
stat(s, Inches(8.8), Inches(4.25), Inches(3.7), "4,672", "moves out", BLUE, 36)
note(s, "Representation is the part you cannot fix later with more compute.")

# =====================================================================
# 3 - Architecture
# =====================================================================
s = slide()
header(s, "The model", "A residual policy–value network")
bullets(s, [
    "3×3 convolution: 18 planes → 64 channels",
    "10 residual blocks, 64 filters, batch normalisation",
    ("Policy head → 4,672 logits   ·   Value head → one scalar", BLUE, True),
    "Both heads share the trunk — the representation must serve two "
    "tasks at once.",
    "Trained with the AlphaZero loss: cross-entropy on the search "
    "distribution, plus MSE on the game result.",
], top=Inches(2.15), width=Inches(7.3), size=18)

stat(s, Inches(8.7), Inches(2.4), Inches(3.8), "~1000×", "smaller than AlphaZero", AMBER)
stat(s, Inches(8.7), Inches(4.25), Inches(3.8), "6 CPU + 1 GPU",
     "self-play on cores, training on GPU", MUTED, 26)
note(s, "Self-play is CPU-bound; the GPU only does batched gradient steps.")

# =====================================================================
# 4 - Stage 1 results
# =====================================================================
s = slide()
header(s, "Stage 1 · supervised", "Learning from 677k human positions",
       GREEN)
picture(s, "pretrain_loss.png", MARGIN, Inches(2.2), Inches(6.0))
bullets(s, [
    ("Lichess database, both players rated 1750+", INK, True),
    "The Elo filter keeps only ~12% of games — imitate competent play, "
    "not average play.",
    ("Policy cross-entropy 3.90 → 2.06", GREEN, True),
    "≈ 13% probability on the exact move a strong human chose, out of "
    "4,672 options.",
    ("It plays real chess: e4/d4 openings, wins a hanging queen, finds mate "
     "in one. 85% vs random.", GREEN, True),
], top=Inches(2.3), left=Inches(7.1), width=Inches(5.5), size=16)

# =====================================================================
# 5 - The loss curve that lied
# =====================================================================
s = slide()
header(s, "Stage 2 · self-play", "300 iterations. The loss fell the whole way.")
picture(s, "v5_regression.png", MARGIN, Inches(2.2), Inches(6.0))
bullets(s, [
    ("Training loss: 2.36 → 1.85", INK, True),
    "Monotonic. Textbook. Exactly what you hope to see.",
    ("Reported evaluation score: 50%", MUTED, True),
    "For 159 consecutive measurements across two multi-hour runs.",
    ("I read that as “evenly matched.”", RED, True),
    "It was not. It was a metric that could not return anything else.",
], top=Inches(2.3), left=Inches(7.1), width=Inches(5.5), size=17)

# =====================================================================
# 6 - The broken metric
# =====================================================================
s = slide()
header(s, "The bug", "The evaluation was measuring nothing", RED)
bullets(s, [
    "The network alternated colours — half its games as white, half as "
    "black.",
    ("The score counted white’s wins.", RED, True),
    "So every game the network won as black was recorded as a loss.",
    "With two deterministic players and no opening randomisation, all five "
    "games per colour are identical — the tally can only be symmetric.",
], top=Inches(2.1), width=Inches(6.55), size=18)

code(s, [
    ("[iter 454] eval: {'1-0': 0, '0-1': 0, '1/2-1/2': 10}  50%", MUTED),
    ("[iter 499] eval: {'1-0': 0, '0-1': 0, '1/2-1/2': 10}  50%", MUTED),
], Inches(7.6), Inches(2.4), Inches(5.0), size=13)

stat(s, Inches(7.6), Inches(4.0), Inches(2.35), "159", "useless data points", RED, 40)
stat(s, Inches(10.25), Inches(4.0), Inches(2.35), "44% → 85%",
     "true score, once fixed", GREEN, 26)
note(s, "A metric returning a plausible constant is more dangerous than one "
        "that crashes.")

# =====================================================================
# 7 - The real result
# =====================================================================
s = slide()
header(s, "The real result", "Self-play made the model worse", RED)

panel(s, MARGIN, Inches(2.3), BODY_W, Inches(1.9), RGBColor(0xFB, 0xF0, 0xEF))
tf = textbox(s, MARGIN, Inches(2.55), BODY_W, Inches(1.5))
para(tf, "v5 final   2.5  —  5.5   its own starting point", 40, RED,
     bold=True, align=PP_ALIGN.CENTER, space_after=6, first=True)
para(tf, "8 games, alternating colours, randomised openings", 15, MUTED,
     align=PP_ALIGN.CENTER, space_after=0)

bullets(s, [
    ("The final model lost to the model it started from.", INK, True),
    "Hours of compute made it measurably weaker — while the loss curve "
    "fell the entire time.",
    ("Why the loss is not a strength signal:", BLUE, True),
    "The loss is computed against targets the model generated itself. "
    "Minimising it proves the model agrees with itself. It says nothing "
    "about whether it plays better chess.",
    "Every self-training and self-supervised pipeline shares this structure "
    "— and therefore this failure mode.",
], top=Inches(4.5), size=17)

# =====================================================================
# 8 - Diagnosis
# =====================================================================
s = slide()
header(s, "Diagnosis", "Two root causes, isolated by experiment")
picture(s, "sims_scaling.png", MARGIN, Inches(2.25), Inches(5.6))

bullets(s, [
    ("1 · The targets were noise", BLUE, True),
    "Search is not weak — MCTS beats the raw policy 79% at 40 "
    "simulations. But 40 simulations over ~30 legal moves is "
    "1.3 visits per move.",
    "The best move is good; the distribution is sampling noise — and "
    "the loss trains on the distribution.",
    ("2 · The augmentation was illegal", RED, True),
    "Mirroring flipped board files but never swapped the castling planes. "
    "The king lands on d1 with “kingside castling” still set — "
    "a position that cannot exist in chess.",
    ("~50% of every training batch was corrupted.", RED, True),
], top=Inches(2.2), left=Inches(6.5), width=Inches(6.1), size=15)

# =====================================================================
# 9 - The fix
# =====================================================================
s = slide()
header(s, "The fix", "Gating stopped the bleeding", GREEN)
bullets(s, [
    ("Acceptance gating, from AlphaGo Zero", GREEN, True),
    "Self-play always generates from the best weights so far; new weights "
    "are promoted only after scoring ≥ 55% against the incumbent.",
    "Plus castling-aware augmentation, 128 simulations instead of 40, "
    "80 gradient steps instead of 200.",
    ("199 iterations later — no regression.", GREEN, True),
    "The candidates the gate rejected averaged 42.2% vs baseline "
    "(p = 0.004): self-play was still degrading the network.",
    "The gate kept that damage out of the deployed model.",
], top=Inches(2.1), width=Inches(7.0), size=16)

code(s, [
    ("gate: candidate 43.8%", INK),
    ("      vs best -> rejected", RED),
], Inches(8.2), Inches(2.4), Inches(4.4), size=15)

stat(s, Inches(8.2), Inches(4.0), Inches(4.4), "46.2%",
     "gated model vs its starting point — no detectable change", GREEN, 38)
note(s, "v5 without the fixes: 31.2%. Gating converted a real regression "
        "into no change.")

# =====================================================================
# 10 - Takeaways
# =====================================================================
s = slide()
header(s, "What I take away", "Measurement is the hard part")
bullets(s, [
    ("A falling loss curve is not evidence that anything is working.", INK, True),
    "Especially when the model produces its own labels.",
    ("Verify the metric before trusting the model.", INK, True),
    "Mine returned a plausible number 159 times in a row while measuring "
    "nothing at all.",
    ("Downscaling changes which statistic matters.", INK, True),
    "800 → 40 simulations preserves the best move but destroys the "
    "distribution — and the loss reads the distribution.",
    ("Then audit the fix too.", RED, True),
    "My gate promoted 8 of 20 candidates. Pure chance predicts 7.8 — "
    "an 8-game match could never have resolved this. The same mistake, "
    "one level up.",
], top=Inches(2.05), width=Inches(7.6), size=16)

panel(s, Inches(8.9), Inches(2.3), Inches(3.6), Inches(3.4))
tf = textbox(s, Inches(9.15), Inches(2.55), Inches(3.1), Inches(3.0))
para(tf, "Artifacts", 17, BLUE, bold=True, space_after=12, first=True)
para(tf, "Code", 13, MUTED, bold=True, space_after=2)
para(tf, "github.com/realgauravvyas/chess-ai", 12, INK, space_after=12)
para(tf, "Results", 13, MUTED, bold=True, space_after=2)
para(tf, "RESULTS.md — every number, reproducible", 12, INK, space_after=12)
para(tf, "Demo", 13, MUTED, bold=True, space_after=2)
para(tf, "Live training dashboard + play any checkpoint", 12, INK, space_after=0)

prs.save(str(OUT))
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
